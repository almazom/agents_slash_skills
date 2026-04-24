#!/usr/bin/env python3
"""
PPTX Post-Translation Overflow Fixer
Fixes layout breakage after translating PPTX files (especially EN→RU, ZH→RU).

Usage:
  python3 pptx_overflow_fix.py input.pptx output.pptx
  python3 pptx_overflow_fix.py input.pptx output.pptx --min-size 8 --lang ru
  python3 pptx_overflow_fix.py input.pptx --report   # Only report, don't fix
"""

import sys
import math
import argparse
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.text import MSO_AUTO_SIZE
    from pptx.util import Pt, Emu
except ImportError:
    print("ERROR: python-pptx required. Install: pip install python-pptx")
    sys.exit(1)

# Text expansion ratios by target language
EXPANSION_RATIOS = {
    "ru": 1.30,  # EN→RU average +30%
    "zh-ru": 1.75,  # ZH→RU average +75%
    "de": 1.30,
    "fr": 1.20,
    "es": 1.20,
    "pt": 1.25,
    "it": 1.20,
    "pl": 1.30,
    "cs": 1.30,
    "fi": 1.35,
    "hu": 1.35,
    "nl": 1.30,
    "ar": 1.30,
    "ja": 0.85,  # Contraction
    "ko": 0.85,
    "zh": 0.75,
}

MIN_FONT_SIZE_PT = 8
DEFAULT_MAX_FONT_PT = 14


def estimate_text_height(text, cx_emu, font_size_pt, margins):
    """Estimate rendered text height in points."""
    avail_w_pt = (cx_emu - margins["l"] - margins["r"]) / 12700
    if avail_w_pt <= 0:
        return 0
    chars_per_line = avail_w_pt / (font_size_pt * 0.52)
    lines = max(1, math.ceil(len(text) / max(chars_per_line, 1)))
    lines += text.count("\n") + text.count("\x0b")  # line breaks
    return lines * font_size_pt * 1.2


def get_shape_margins(tf):
    """Get text frame margins in EMU."""
    return {
        "l": tf.margin_left if tf.margin_left else 91440,
        "r": tf.margin_right if tf.margin_right else 91440,
        "t": tf.margin_top if tf.margin_top else 45720,
        "b": tf.margin_bottom if tf.margin_bottom else 45720,
    }


def get_font_size(tf):
    """Get the first non-None font size from text frame, in EMU."""
    for para in tf.paragraphs:
        for run in para.runs:
            if run.font.size:
                return run.font.size
    return 183000  # Default 14.4pt


def detect_overflow(prs):
    """Detect all text frames with overflow."""
    results = []
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text
            if not text.strip():
                continue

            margins = get_shape_margins(tf)
            cx = shape.width
            cy = shape.height
            font_size_emu = get_font_size(tf)
            font_size_pt = font_size_emu / 12700

            effective_h = (cy - margins["t"] - margins["b"]) / 12700
            estimated_h = estimate_text_height(text, cx, font_size_pt, margins)

            if estimated_h > effective_h and effective_h > 0:
                ratio = estimated_h / effective_h
                results.append({
                    "slide": slide_num,
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name if hasattr(shape, "name") else "unknown",
                    "text_preview": text[:80].replace("\n", " "),
                    "font_pt": round(font_size_pt, 1),
                    "overflow_ratio": round(ratio, 2),
                    "severity": "severe" if ratio > 1.8 else "moderate" if ratio > 1.3 else "mild",
                })
    return results


def fix_overflows(prs, min_size=MIN_FONT_SIZE_PT, max_size=DEFAULT_MAX_FONT_PT):
    """Apply tiered overflow fixes to all text frames."""
    fixes_applied = {"mild": 0, "moderate": 0, "severe": 0}

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            tf = shape.text_frame
            text = tf.text
            if not text.strip():
                continue

            margins = get_shape_margins(tf)
            cx = shape.width
            cy = shape.height
            font_size_emu = get_font_size(tf)
            font_size_pt = font_size_emu / 12700
            effective_h = (cy - margins["t"] - margins["b"]) / 12700
            estimated_h = estimate_text_height(text, cx, font_size_pt, margins)

            if estimated_h <= effective_h or effective_h <= 0:
                continue

            ratio = estimated_h / effective_h

            if ratio <= 1.3:
                # Tier 1: Mild — set normAutofit
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                fixes_applied["mild"] += 1

            elif ratio <= 1.8:
                # Tier 2: Moderate — shrink font + normAutofit
                try:
                    tf.fit_text(
                        font_family="Arial",
                        max_size=min(int(font_size_pt), max_size),
                    )
                except Exception:
                    # fit_text can fail on edge cases; fallback to manual shrink
                    new_size = max(min_size, int(font_size_pt / ratio))
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                run.font.size = Pt(new_size)
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                fixes_applied["moderate"] += 1

            else:
                # Tier 3: Severe — reduce margins + shrink + expand box
                tf.margin_left = tf.margin_right = 45720
                tf.margin_top = tf.margin_bottom = 22860
                try:
                    tf.fit_text(font_family="Arial", max_size=min(int(font_size_pt), max_size))
                except Exception:
                    new_size = max(min_size, int(font_size_pt * 0.7))
                    for para in tf.paragraphs:
                        for run in para.runs:
                            if run.font.size:
                                run.font.size = Pt(new_size)
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                # Expand shape height by expansion ratio
                new_height = int(cy * min(ratio, 1.5))
                shape.height = new_height
                fixes_applied["severe"] += 1

    return fixes_applied


def main():
    parser = argparse.ArgumentParser(description="PPTX Post-Translation Overflow Fixer")
    parser.add_argument("input", help="Input PPTX file")
    parser.add_argument("output", nargs="?", help="Output PPTX file (omit for report-only)")
    parser.add_argument("--min-size", type=int, default=MIN_FONT_SIZE_PT, help="Minimum font size in pt")
    parser.add_argument("--max-size", type=int, default=DEFAULT_MAX_FONT_PT, help="Maximum font size in pt")
    parser.add_argument("--report", action="store_true", help="Only report overflows, don't fix")
    parser.add_argument("--lang", default="ru", help="Target language code for expansion ratio reference")

    args = parser.parse_args()

    if not Path(args.input).exists():
        print(f"ERROR: File not found: {args.input}")
        sys.exit(1)

    prs = Presentation(args.input)
    total_slides = len(prs.slides)

    print(f"Analyzing: {args.input}")
    print(f"Slides: {total_slides}")
    print(f"Expected expansion ratio for {args.lang}: {EXPANSION_RATIOS.get(args.lang, 1.3):.2f}x")
    print()

    overflows = detect_overflow(prs)

    if not overflows:
        print("No text overflow detected. Document looks clean.")
        return

    print(f"OVERFLOW DETECTED: {len(overflows)} shapes")
    print("=" * 70)
    for o in overflows:
        print(f"  Slide {o['slide']:>2} | {o['severity']:>8} | ratio {o['overflow_ratio']:.2f} | "
              f"{o['font_pt']}pt | {o['text_preview'][:50]}")
    print()

    if args.report:
        print("Report mode — no fixes applied.")
        return

    if not args.output:
        print("No output file specified. Use --report for analysis only, or provide output path.")
        sys.exit(1)

    # Apply fixes
    fixes = fix_overflows(prs, min_size=args.min_size, max_size=args.max_size)

    # Write to /tmp first (avoid Docker bind-mount seek issues)
    tmp_out = Path(f"/tmp/{Path(args.output).name}")
    prs.save(str(tmp_out))

    # Copy to final destination
    import shutil
    shutil.copy2(str(tmp_out), args.output)

    print(f"FIXES APPLIED:")
    print(f"  Mild (normAutofit):     {fixes['mild']}")
    print(f"  Moderate (shrink font): {fixes['moderate']}")
    print(f"  Severe (resize box):    {fixes['severe']}")
    print(f"\nSaved: {args.output}")


if __name__ == "__main__":
    main()
